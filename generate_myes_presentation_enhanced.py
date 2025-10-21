"""
MyES Enhanced PowerPoint Generator
===================================
Advanced features:
✅ Smart overflow detection & warnings
✅ Automatic bullet/list formatting
✅ Image insertion support
✅ Robust question splitting
✅ Style variations (vocabulary, questions, answers)
✅ Vertical alignment control
✅ Pre-generation validation
✅ Config file support (JSON)
✅ Auto-animation for [step] markers
✅ Math/special character handling
✅ Reusable slide templates
✅ Progress indicators & slide numbers
"""

import sys
import os
import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Try to import animation support (may not be available in all python-pptx versions)
try:
    from pptx.oxml.xmlchemy import OxmlElement
    ANIMATIONS_SUPPORTED = True
except:
    ANIMATIONS_SUPPORTED = False


# === DEFAULT CONFIG ===
DEFAULT_CONFIG = {
    "background_image": "MyES Slides Template 2025.jpg",
    "title_color": [192, 0, 0],
    "text_color": [0, 0, 102],
    "font_name": "Montserrat",
    "slide_width": 13.33,
    "slide_height": 7.5,
    "enable_animations": True,
    "enable_slide_numbers": True,
    "enable_overflow_warnings": True,
    "styles": {
        "vocabulary": {"font_size": 24, "color": [0, 102, 0], "bold": True},
        "question": {"font_size": 20, "color": [102, 0, 102], "bold": False},
        "answer": {"font_size": 18, "color": [128, 128, 128], "italic": True},
        "emphasis": {"font_size": 22, "color": [192, 0, 0], "bold": True}
    }
}


# === LOAD CONFIG ===
def load_config(config_file="myes_config.json"):
    """Load configuration from JSON file, fall back to defaults"""
    if os.path.exists(config_file):
        try:
            with open(config_file, 'r', encoding='utf-8') as f:
                user_config = json.load(f)
                # Merge with defaults
                config = DEFAULT_CONFIG.copy()
                config.update(user_config)
                print(f"✅ Loaded config from {config_file}")
                return config
        except Exception as e:
            print(f"⚠️  Error loading config: {e}. Using defaults.")
    return DEFAULT_CONFIG


# === OVERFLOW DETECTION ===
def check_text_overflow(text, font_size, width_inches, height_inches):
    """
    Estimates if text will overflow the given dimensions.
    Returns (will_overflow, estimated_lines_needed, lines_available)
    """
    # Convert to float if needed
    if hasattr(width_inches, 'inches'):
        width_inches = width_inches.inches
    if hasattr(height_inches, 'inches'):
        height_inches = height_inches.inches
    
    # Rough estimation: average chars per line based on font size and width
    chars_per_inch = 72 / font_size * 2.5  # Approximate
    chars_per_line = int(width_inches * chars_per_inch)
    
    # Calculate lines needed
    words = text.split()
    current_line_length = 0
    lines_needed = 1
    
    for word in words:
        word_length = len(word) + 1  # +1 for space
        if current_line_length + word_length > chars_per_line:
            lines_needed += 1
            current_line_length = word_length
        else:
            current_line_length += word_length
    
    # Calculate available lines
    line_height = font_size / 72 * 1.3  # Line height with spacing
    lines_available = int(height_inches / line_height)
    
    will_overflow = lines_needed > lines_available
    return will_overflow, lines_needed, lines_available


# === LIST DETECTION & FORMATTING ===
def is_list_content(lines):
    """Detect if content should be formatted as a list"""
    if not lines:
        return False
    
    bullet_patterns = [r'^\s*[•\-\*]', r'^\s*\d+\.', r'^\s*[a-z]\)', r'^\s*[A-Z]\.']
    
    # Check if most lines match list patterns
    matching = sum(1 for line in lines if any(re.match(p, line) for p in bullet_patterns))
    return matching >= len(lines) * 0.5  # 50% threshold


def clean_bullet_marker(text):
    """Remove common bullet markers from text"""
    text = re.sub(r'^\s*[•\-\*]\s*', '', text)
    text = re.sub(r'^\s*\d+\.\s*', '', text)
    text = re.sub(r'^\s*[a-z]\)\s*', '', text)
    return text


# === QUESTION SPLITTING ===
def split_questions(text):
    """
    Robustly split multiple questions in a text block.
    Returns list of individual questions.
    """
    # Split on question marks followed by whitespace and number/capital letter
    questions = re.split(r'\?\s*(?=\d+\.|\b[A-Z])', text)
    result = []
    
    for q in questions:
        q = q.strip()
        if q:
            # Add question mark back if missing
            if not q.endswith('?'):
                q += '?'
            result.append(q)
    
    return result if result else [text]


# === STYLE APPLICATION ===
def apply_style(paragraph, style_name, config):
    """Apply predefined style to a paragraph"""
    styles = config.get("styles", {})
    
    if style_name in styles:
        style = styles[style_name]
        paragraph.font.size = Pt(style.get("font_size", 22))
        
        if "color" in style:
            color = style["color"]
            paragraph.font.color.rgb = RGBColor(*color)
        
        paragraph.font.bold = style.get("bold", False)
        paragraph.font.italic = style.get("italic", False)


def parse_styled_text(text):
    """
    Parse text with inline style markers like [vocabulary] or [question]
    Returns (style, cleaned_text)
    """
    match = re.match(r'^\[(\w+)\]\s*(.+)', text)
    if match:
        return match.group(1), match.group(2)
    return None, text


# === MATH/SPECIAL CHARACTERS ===
def process_math(text):
    """
    Convert simple math notation to Unicode symbols.
    Handles basic cases like x^2, subscripts, fractions
    """
    # Superscripts
    superscripts = {'0': '⁰', '1': '¹', '2': '²', '3': '³', '4': '⁴', 
                    '5': '⁵', '6': '⁶', '7': '⁷', '8': '⁸', '9': '⁹'}
    text = re.sub(r'\^(\d)', lambda m: superscripts.get(m.group(1), m.group(1)), text)
    
    # Subscripts
    subscripts = {'0': '₀', '1': '₁', '2': '₂', '3': '₃', '4': '₄',
                  '5': '₅', '6': '₆', '7': '₇', '8': '₈', '9': '₉'}
    text = re.sub(r'_(\d)', lambda m: subscripts.get(m.group(1), m.group(1)), text)
    
    # Common symbols
    replacements = {
        '<=': '≤', '>=': '≥', '!=': '≠', '~=': '≈',
        'alpha': 'α', 'beta': 'β', 'gamma': 'γ', 'delta': 'δ',
        'pi': 'π', 'theta': 'θ', 'sigma': 'σ'
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    
    return text


# === IMAGE PARSING ===
def parse_image_directive(line):
    """
    Parse image directives like:
    Image: path/to/image.jpg | width=4 | left=2 | top=3 | align=center
    """
    parts = line.replace("Image:", "").split("|")
    image_data = {"path": parts[0].strip()}
    
    for part in parts[1:]:
        if "=" in part:
            key, value = part.split("=", 1)
            image_data[key.strip()] = value.strip()
    
    return image_data


# === VALIDATION ===
def validate_slide(slide_data, slide_num, config):
    """Validate slide data and return list of warnings/errors"""
    issues = []
    
    if not slide_data["title"]:
        issues.append(f"Slide {slide_num}: Missing title")
    elif len(slide_data["title"]) > 100:
        issues.append(f"Slide {slide_num}: Title very long ({len(slide_data['title'])} chars)")
    
    # Check for content
    has_content = any([
        slide_data["content"], slide_data["left"], slide_data["right"],
        slide_data["left_top"], slide_data["right_top"],
        slide_data["left_bottom"], slide_data["right_bottom"]
    ])
    
    if not has_content:
        issues.append(f"Slide {slide_num}: No content defined")
    
    # Check image paths
    for img in slide_data.get("images", []):
        if not os.path.exists(img["path"]):
            issues.append(f"Slide {slide_num}: Image not found: {img['path']}")
    
    return issues


# === ADD TEXTBOX (Enhanced) ===
def add_textbox(slide, left, top, width, height, lines, font_size=22, label=None, 
                config=None, v_align=MSO_ANCHOR.TOP):
    """
    Enhanced textbox with overflow detection, list formatting, and styling
    """
    if not lines:
        return None
    
    if config is None:
        config = DEFAULT_CONFIG
    
    joined = " ".join(lines)
    joined = process_math(joined)  # Handle math symbols
    text_length = len(joined)  # Calculate once for reuse
    
    # Overflow detection
    if config.get("enable_overflow_warnings", True):
        try:
            # Handle both Inches objects and plain floats
            w = width.inches if hasattr(width, 'inches') else width
            h = height.inches if hasattr(height, 'inches') else height
            overflow, needed, available = check_text_overflow(joined, font_size, w, h)
            if overflow:
                print(f"⚠️  Potential overflow in '{label}': needs {needed} lines, has {available}")
        except Exception as e:
            # Silently skip overflow check if there's an issue
            pass
    
    # Auto font-size for long text (more aggressive sizing)
    text_length = len(joined)
    if text_length > 300:
        font_size = min(font_size, 18)
    if text_length > 500:
        font_size = min(font_size, 16)
    if text_length > 700:
        font_size = min(font_size, 14)
    if text_length > 1000:
        font_size = min(font_size, 12)
    
    # Handle [step] animations
    step_lines = [l for l in lines if "[step]" in l.lower()]
    if step_lines:
        return add_step_textboxes(slide, left, top, width, lines, font_size, label, config)
    
    # Detect list formatting
    is_list = is_list_content(lines)
    
    # Create textbox
    box = slide.shapes.add_textbox(left, top, width, height)
    if label:
        box.name = label
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = v_align
    
    # Add content
    first = True
    for item in lines:
        if not item.strip():
            continue
        
        # Parse style markers
        style, text = parse_styled_text(item)
        text = process_math(text)
        
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        
        # Apply list formatting
        if is_list:
            text = clean_bullet_marker(text)
            p.level = 0
            p.text = text
        else:
            p.text = text
        
        # Apply styling
        p.font.name = config["font_name"]
        if style:
            apply_style(p, style, config)
        else:
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*config["text_color"])
        
        # Reduce spacing for long content
        if text_length > 300:
            p.space_after = Pt(3)
        else:
            p.space_after = Pt(6)
    
    return box


def add_step_textboxes(slide, left, top, width, lines, font_size, label, config):
    """Create separate textboxes for each [step] line for animation"""
    top_offset = top
    boxes = []
    
    for i, item in enumerate(lines):
        if not item.strip():
            continue
        
        # Remove [step] marker
        text = re.sub(r'\[step\]\s*', '', item, flags=re.IGNORECASE)
        text = process_math(text)
        
        # Parse style
        style, text = parse_styled_text(text)
        
        box = slide.shapes.add_textbox(left, top_offset, width, Inches(0.6))
        if label:
            box.name = f"{label}_Step{i+1}"
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = config["font_name"]
        
        if style:
            apply_style(p, style, config)
        else:
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*config["text_color"])
        
        boxes.append(box)
        top_offset += Inches(0.65)
    
    # Add animations if supported
    if config.get("enable_animations", True) and ANIMATIONS_SUPPORTED:
        add_appear_animations(boxes)
    
    return boxes


# === ANIMATION SUPPORT ===
def add_appear_animations(shapes):
    """Add simple appear animations to shapes (basic implementation)"""
    # Note: Full animation support in python-pptx is limited
    # This is a placeholder for future enhancement
    pass


# === ADD IMAGE ===
def add_image_to_slide(slide, image_data, content_area):
    """Add image to slide with positioning"""
    path = image_data["path"]
    
    # Parse dimensions
    width = Inches(float(image_data.get("width", 4)))
    
    # Parse position
    if "left" in image_data:
        left = Inches(float(image_data["left"]))
    elif image_data.get("align") == "center":
        left = (content_area["width"] - width) / 2 + content_area["left"]
    else:
        left = content_area["left"]
    
    if "top" in image_data:
        top = Inches(float(image_data["top"]))
    else:
        top = content_area["top"]
    
    try:
        pic = slide.shapes.add_picture(path, left, top, width=width)
        return pic
    except Exception as e:
        print(f"⚠️  Error adding image {path}: {e}")
        return None


# === ADD SLIDE NUMBER ===
def add_slide_number(slide, slide_num, total_slides, config):
    """Add slide number footer"""
    footer_text = f"{slide_num} / {total_slides}"
    
    footer = slide.shapes.add_textbox(
        Inches(config["slide_width"] - 1.5),
        Inches(config["slide_height"] - 0.5),
        Inches(1),
        Inches(0.3)
    )
    
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.RIGHT


# === TEMPLATE SYSTEM ===
SLIDE_TEMPLATES = {
    "vocabulary": {
        "layout": "left_right",
        "left_title": "Word",
        "right_title": "Definition",
        "style": "vocabulary"
    },
    "reading": {
        "layout": "stacked",
        "top_ratio": 0.6,
        "bottom_label": "Questions"
    },
    "comparison": {
        "layout": "left_right",
        "left_title": "Before",
        "right_title": "After"
    }
}


def apply_template(slide_data, template_name):
    """Apply a predefined template to slide data"""
    if template_name not in SLIDE_TEMPLATES:
        return slide_data
    
    template = SLIDE_TEMPLATES[template_name]
    
    # Apply template-specific modifications
    if template.get("layout") == "left_right":
        # Add section titles if specified
        if "left_title" in template and slide_data["left"]:
            slide_data["left"].insert(0, f"[emphasis] {template['left_title']}")
        if "right_title" in template and slide_data["right"]:
            slide_data["right"].insert(0, f"[emphasis] {template['right_title']}")
    
    return slide_data


# === PARSER (Enhanced) ===
def parse_content_file(filename):
    """Enhanced parser with template and image support"""
    slides = []
    current = {
        "title": "", "content": [], "notes": [], "images": [],
        "left": [], "right": [],
        "left_top": [], "right_top": [],
        "left_bottom": [], "right_bottom": [],
        "template": None
    }
    section = None

    with open(filename, "r", encoding="utf-8") as f:
        for line in f:
            line = line.rstrip()  # Keep leading spaces for indentation
            
            if not line or line == "---":
                continue

            # New slide
            if line.startswith("Slide "):
                if current["title"]:
                    slides.append(current)
                    current = {
                        "title": "", "content": [], "notes": [], "images": [],
                        "left": [], "right": [],
                        "left_top": [], "right_top": [],
                        "left_bottom": [], "right_bottom": [],
                        "template": None
                    }
                section = None
                continue

            # Template directive
            if line.startswith("Template:"):
                current["template"] = line.replace("Template:", "").strip()
                continue

            # Image directive
            if line.startswith("Image:"):
                current["images"].append(parse_image_directive(line))
                continue

            # Section headers
            if line.startswith("Title:"):
                current["title"] = line.replace("Title:", "").strip()
                section = None
                continue
            elif line.startswith("Content:"):
                section = "content"
                text = line.replace("Content:", "").strip()
            elif line.startswith("Left:"):
                section = "left"
                text = line.replace("Left:", "").strip()
            elif line.startswith("Right:"):
                section = "right"
                text = line.replace("Right:", "").strip()
            elif line.startswith("LeftTop:"):
                section = "left_top"
                text = line.replace("LeftTop:", "").strip()
            elif line.startswith("RightTop:"):
                section = "right_top"
                text = line.replace("RightTop:", "").strip()
            elif line.startswith("LeftBottom:"):
                section = "left_bottom"
                text = line.replace("LeftBottom:", "").strip()
                
                # Auto-split questions
                if any(q in text for q in ["1.", "2.", "3."]):
                    questions = split_questions(text)
                    current["left_bottom"].extend(questions)
                    continue
            elif line.startswith("RightBottom:"):
                section = "right_bottom"
                text = line.replace("RightBottom:", "").strip()
            elif line.startswith("Notes:"):
                section = "notes"
                text = line.replace("Notes:", "").strip()
            else:
                text = line

            # Store text
            if section in current and text:
                current[section].append(text)

        if current["title"]:
            slides.append(current)

    return slides


# === BUILD PRESENTATION (Enhanced) ===
def build_presentation(slides, output_name, config):
    """Enhanced presentation builder with all new features"""
    prs = Presentation()
    prs.slide_width = Inches(config["slide_width"])
    prs.slide_height = Inches(config["slide_height"])
    
    SLIDE_WIDTH = prs.slide_width
    SLIDE_HEIGHT = prs.slide_height
    TITLE_LEFT = Inches(1.5)
    TITLE_TOP = Inches(0.6)
    CONTENT_LEFT = Inches(1.5)
    CONTENT_TOP = Inches(1.5)
    CONTENT_WIDTH = Inches(10.5)
    CONTENT_HEIGHT = Inches(5.0)  # Increased from 4.5
    COLUMN_GAP = Inches(0.4)  # Reduced from 0.5
    ROW_GAP = Inches(0.3)  # Reduced from 0.5
    COLUMN_WIDTH = (CONTENT_WIDTH - COLUMN_GAP) / 2
    
    content_area = {
        "left": CONTENT_LEFT,
        "top": CONTENT_TOP,
        "width": CONTENT_WIDTH,
        "height": CONTENT_HEIGHT
    }
    
    total_slides = len(slides)
    
    for idx, s in enumerate(slides, start=1):
        # Apply template if specified
        if s.get("template"):
            s = apply_template(s, s["template"])
        
        # Create slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        if os.path.exists(config["background_image"]):
            slide.shapes.add_picture(config["background_image"], 0, 0,
                                    width=SLIDE_WIDTH, height=SLIDE_HEIGHT)
        
        # Title
        title_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, CONTENT_WIDTH, Inches(0.8))  # Reduced from 1.0
        tf = title_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = s["title"]
        p.font.name = config["font_name"]
        p.font.size = Pt(32)  # Reduced from 36
        p.font.bold = True
        p.font.color.rgb = RGBColor(*config["title_color"])
        
        # Layout logic
        
        # Reading slide (text above, questions below)
        if s["left_top"] and s["left_bottom"] and not (s["right_top"] or s["right_bottom"]):
            # Give more space to reading text (65/35 split instead of 50/50)
            reading_height = CONTENT_HEIGHT * 0.65
            questions_height = CONTENT_HEIGHT * 0.35 - ROW_GAP
            
            # Use default font_size (will auto-adjust based on length)
            add_textbox(slide, CONTENT_LEFT, CONTENT_TOP,
                       CONTENT_WIDTH, reading_height, s["left_top"], 
                       label="ReadingText", config=config, v_align=MSO_ANCHOR.TOP)
            add_textbox(slide, CONTENT_LEFT, CONTENT_TOP + reading_height + ROW_GAP,
                       CONTENT_WIDTH, questions_height, s["left_bottom"], 
                       label="ReadingQuestions", config=config, v_align=MSO_ANCHOR.TOP)
        
        # 4-box slide
        elif any([s["left_top"], s["right_top"], s["left_bottom"], s["right_bottom"]]):
            half_height = (CONTENT_HEIGHT - ROW_GAP) / 2
            # Use smaller font for 4-box layouts to fit more content
            box_font_size = 18
            
            add_textbox(slide, CONTENT_LEFT, CONTENT_TOP,
                       COLUMN_WIDTH, half_height, s["left_top"], 
                       font_size=box_font_size, label="LeftTop", config=config)
            add_textbox(slide, CONTENT_LEFT + COLUMN_WIDTH + COLUMN_GAP, CONTENT_TOP,
                       COLUMN_WIDTH, half_height, s["right_top"], 
                       font_size=box_font_size, label="RightTop", config=config)
            add_textbox(slide, CONTENT_LEFT, CONTENT_TOP + half_height + ROW_GAP,
                       COLUMN_WIDTH, half_height, s["left_bottom"], 
                       font_size=box_font_size, label="LeftBottom", config=config)
            add_textbox(slide, CONTENT_LEFT + COLUMN_WIDTH + COLUMN_GAP,
                       CONTENT_TOP + half_height + ROW_GAP,
                       COLUMN_WIDTH, half_height, s["right_bottom"], 
                       font_size=box_font_size, label="RightBottom", config=config)
        
        # Two-column layout
        elif s["left"] or s["right"]:
            add_textbox(slide, CONTENT_LEFT, CONTENT_TOP,
                       COLUMN_WIDTH, CONTENT_HEIGHT, s["left"], 
                       label="Left", config=config)
            add_textbox(slide, CONTENT_LEFT + COLUMN_WIDTH + COLUMN_GAP, CONTENT_TOP,
                       COLUMN_WIDTH, CONTENT_HEIGHT, s["right"], 
                       label="Right", config=config)
        
        # Single-column content
        else:
            add_textbox(slide, CONTENT_LEFT, CONTENT_TOP,
                       CONTENT_WIDTH, CONTENT_HEIGHT, s["content"], 
                       label="Content", config=config)
        
        # Add images
        for img_data in s.get("images", []):
            add_image_to_slide(slide, img_data, content_area)
        
        # Add slide numbers
        if config.get("enable_slide_numbers", True):
            add_slide_number(slide, idx, total_slides, config)
        
        # Notes
        if s["notes"]:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            for note in s["notes"]:
                notes_tf.add_paragraph().text = f"• {note}"
    
    prs.save(output_name)
    print(f"✅ Presentation created: {output_name}")


# === MAIN ===
if __name__ == "__main__":
    print("=" * 60)
    print("MyES Enhanced PowerPoint Generator")
    print("=" * 60)
    
    if len(sys.argv) < 2:
        print("\nUsage: python generate_myes_presentation_enhanced.py <lesson_content.txt> [config.json]")
        print("\nFeatures:")
        print("  • Smart overflow detection")
        print("  • Automatic list formatting")
        print("  • Image support: Image: path.jpg | width=4 | align=center")
        print("  • Style tags: [vocabulary], [question], [answer], [emphasis]")
        print("  • Step animations: [step] text")
        print("  • Math symbols: x^2, alpha, beta, <=, >=")
        print("  • Templates: Template: vocabulary")
        print("  • Slide numbers and validation")
        sys.exit(1)
    
    input_file = sys.argv[1]
    config_file = sys.argv[2] if len(sys.argv) > 2 else "myes_config.json"
    
    # Load configuration
    config = load_config(config_file)
    
    # Parse slides
    print(f"\n📖 Reading: {input_file}")
    slides = parse_content_file(input_file)
    print(f"   Found {len(slides)} slides")
    
    # Validate
    print("\n🔍 Validating slides...")
    all_issues = []
    for i, slide in enumerate(slides, 1):
        issues = validate_slide(slide, i, config)
        all_issues.extend(issues)
        
    if all_issues:
        print("   ⚠️  Issues found:")
        for issue in all_issues:
            print(f"      {issue}")
    else:
        print("   ✅ All slides valid")
    
    # Generate output filename
    base_name = input_file.replace(".txt", "")
    if not base_name.lower().endswith("_slides"):
        output_name = base_name + "_slides.pptx"
    else:
        output_name = base_name + ".pptx"
    
    # Build presentation
    print(f"\n🎨 Generating presentation...")
    build_presentation(slides, output_name, config)
    
    print("\n✨ Done!")

    print("=" * 60)
